from file_converter.core import convert_pptx_to_pdf
from pptx import Presentation


def test_pptx_to_pdf(tmp_path):
    f = tmp_path / "slides.pptx"
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    prs.slides.add_slide(slide_layout)
    prs.save(f)
    out = convert_pptx_to_pdf(f, tmp_path)
    assert out.exists()

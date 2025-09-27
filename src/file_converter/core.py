import os
from pathlib import Path
from typing import Optional

def validate_input_path(input_path: str | os.PathLike) -> Path:
    """
    Valida se o caminho existe, é um arquivo e tem extensão suportada.
    """
    p = Path(input_path)
    if not p.exists():
        raise FileNotFoundError(f"Caminho não existe: {p}")
    if not p.is_file():
        raise IsADirectoryError(f"Esperado arquivo, mas é diretório: {p}")
    if p.suffix.lower() not in {".docx", ".csv", ".pptx"}:
        raise ValueError("Extensão inválida (use .docx, .csv ou .pptx)")
    return p.resolve()

def build_output_path(input_path: str | os.PathLike,
                      output_dir: Optional[str | os.PathLike] = None,
                      suffix: str = ".pdf") -> Path:
    in_path = Path(input_path)
    out_dir = Path(output_dir) if output_dir else in_path.parent
    out_dir = out_dir.resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    return out_dir / (in_path.stem + suffix)

def convert_docx_to_pdf(input_path, output_dir=None):
    from docx2pdf import convert
    src = validate_input_path(input_path)
    out_path = build_output_path(src, output_dir, ".pdf")
    convert(str(src), str(out_path.parent))
    return out_path

def convert_csv_to_pdf(input_path, output_dir=None):
    import pandas as pd
    from reportlab.pdfgen import canvas
    src = validate_input_path(input_path)
    out_path = build_output_path(src, output_dir, ".pdf")
    df = pd.read_csv(src)
    c = canvas.Canvas(str(out_path))
    text = c.beginText(40, 800)
    text.setFont("Helvetica", 10)
    for line in df.to_string().splitlines():
        text.textLine(line)
    c.drawText(text)
    c.save()
    return out_path

def convert_pptx_to_pdf(input_path, output_dir=None):
    from pptx import Presentation
    from reportlab.pdfgen import canvas
    src = validate_input_path(input_path)
    out_path = build_output_path(src, output_dir, ".pdf")
    prs = Presentation(src)
    c = canvas.Canvas(str(out_path))
    for i, slide in enumerate(prs.slides, start=1):
        text = c.beginText(40, 800)
        text.setFont("Helvetica", 12)
        text.textLine(f"Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.textLine(shape.text)
        c.drawText(text)
        c.showPage()
    c.save()
    return out_path
